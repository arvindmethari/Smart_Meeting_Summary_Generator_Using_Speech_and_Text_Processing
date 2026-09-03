import os
import re
import torch
import whisper
from moviepy import VideoFileClip 
from pypdf import PdfReader       
from docx import Document
from pptx import Presentation
from openai import OpenAI
from dotenv import load_dotenv
from transformers import pipeline
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM
from datasets import Dataset
from tqdm import tqdm
load_dotenv()

os.makedirs(r"D:\huggingface_cache", exist_ok=True)
os.environ['HF_HOME'] = r"D:\huggingface_cache"

# Detect hardware acceleration (GPU, Apple Silicon, or CPU)
device = 0 if torch.cuda.is_available() else (-1 if not torch.backends.mps.is_available() else "mps")

# ==============================================================================
# 1. FILE EXTRACTION & CLEANING MODULE
# ==============================================================================

def clean_repeating_text(text):
    # Removes consecutive repeating words
    text = re.sub(r'\b(\w+)( \1\b)+', r'\1', text, flags=re.IGNORECASE)
    words = text.split()
    window_size = 15
    
    # Removes AI transcription loops (hallucinations)
    if len(words) > window_size:
        for i in range(len(words) - window_size):
            window = words[i:i+window_size]
            unique_words = set(re.sub(r'[^a-zA-Z]', '', w.lower()) for w in window if re.sub(r'[^a-zA-Z]', '', w.lower()))
            
            if len(unique_words) <= 5:
                clean_part = " ".join(words[:i])
                last_punctuation = max(clean_part.rfind('.'), clean_part.rfind('?'), clean_part.rfind('!'))
                if last_punctuation != -1:
                    text = clean_part[:last_punctuation + 1]
                else:
                    text = clean_part
                break

    # Reconstruct sentences
    sentences = re.split(r'(?<=[.!?]) +', text)
    cleaned_sentences = []
    for sentence in sentences:
        sentence = sentence.strip()
        if not sentence: continue
        if not cleaned_sentences or cleaned_sentences[-1].lower() != sentence.lower():
            cleaned_sentences.append(sentence)
            
    return " ".join(cleaned_sentences)

def extract_text_from_audio(audio_path, model_size="base"):
    print(f"Transcribing audio: {audio_path}")
    model = whisper.load_model(model_size)
    result = model.transcribe(
        audio_path,
        language="en",
        condition_on_previous_text=False, 
        no_speech_threshold=0.6,          
        logprob_threshold=-1.0,           
        compression_ratio_threshold=2.0   
    )

    return clean_repeating_text(result["text"])


def extract_text_from_video(video_path):
    print(f"Extracting audio from video: {video_path}")
    video = VideoFileClip(video_path)
    # Generate a unique temp filename so it doesn't overwrite existing files
    import uuid
    audio_path = f"temp_audio_{uuid.uuid4().hex}.wav"
    video.audio.write_audiofile(audio_path, logger=None, ffmpeg_params=["-ac", "1"])
    video.close()
    
    text = extract_text_from_audio(audio_path)
    if os.path.exists(audio_path):
        os.remove(audio_path)
    return text

def extract_text_from_document(doc_path):
    ext = doc_path.split('.')[-1].lower()
    text = ""
    try:
        if ext == 'pdf':
            reader = PdfReader(doc_path)
            for page in reader.pages: text += (page.extract_text() or "") + "\n"
        elif ext == 'docx':
            doc = Document(doc_path)
            for para in doc.paragraphs: text += para.text + "\n"
        elif ext == 'pptx':
            prs = Presentation(doc_path)
            for slide in prs.slides:
                for shape in slide.shapes:
                    if hasattr(shape, "text"): text += shape.text + "\n"
    except Exception as e:
        print(f"Error reading document: {e}")
    return text.strip()

# ==============================================================================
# 3. ZERO-SHOT INFERENCE PIPELINE
# ==============================================================================
import re
import torch
from transformers import AutoTokenizer, AutoModelForSeq2SeqLM, pipeline

import re

def clean_transcript_for_summarizer(text: str) -> str:
    # remove AMI style tags or whisper artifacts if present
    text = re.sub(r"\{.*?\}", " ", text)          # {vocalsound}, {gap}, etc
    text = re.sub(r"\s+", " ", text).strip()
    return text

def split_into_sentence_chunks(text: str, max_chars: int = 3500):
    # simple sentence split (fast, no extra deps)
    sents = re.split(r'(?<=[.!?])\s+', text)
    chunks = []
    cur = ""
    for s in sents:
        if not s:
            continue
        if len(cur) + len(s) + 1 <= max_chars:
            cur = (cur + " " + s).strip()
        else:
            if cur:
                chunks.append(cur)
            cur = s
    if cur:
        chunks.append(cur)
    return chunks


class NLPProcessor:
    def __init__(self):
        self.device = 0 if torch.cuda.is_available() else -1
        print("NLP System Ready (Storage set to D:)")

        self.SUMMARY_MODEL_PATH = r"D:\bart_samsum_best"

        self.sum_tokenizer = AutoTokenizer.from_pretrained(self.SUMMARY_MODEL_PATH)
        self.sum_model = AutoModelForSeq2SeqLM.from_pretrained(self.SUMMARY_MODEL_PATH)

        if torch.cuda.is_available():
            self.sum_model = self.sum_model.to("cuda")
        else:
            self.sum_model = self.sum_model.to("cpu").float()

    # ==========================================================
    # SUMMARY GENERATION (LONG + FULL COVERAGE)
    # ==========================================================

    def generate_summary(self, text, doc_type="meeting"):

        text = clean_transcript_for_summarizer(text)
        chunks = split_into_sentence_chunks(text, max_chars=3000)

        partial_summaries = []
        self.sum_model.eval()

        for chunk in chunks:
            inputs = self.sum_tokenizer(
                chunk,
                return_tensors="pt",
                truncation=True,
                max_length=1024
            )

            if torch.cuda.is_available():
                inputs = {k: v.to("cuda") for k, v in inputs.items()}

            with torch.no_grad():
                output_ids = self.sum_model.generate(
                    **inputs,
                    num_beams=6,
                    max_new_tokens=400,
                    min_new_tokens=200,
                    no_repeat_ngram_size=3,
                    repetition_penalty=1.2,
                    length_penalty=1.2,
                    early_stopping=True
                )

            summary = self.sum_tokenizer.decode(
                output_ids[0],
                skip_special_tokens=True
            )

            partial_summaries.append(summary)

        # Second pass summarization
        combined = " ".join(partial_summaries)

        inputs = self.sum_tokenizer(
            combined,
            return_tensors="pt",
            truncation=True,
            max_length=1024
        )

        if torch.cuda.is_available():
            inputs = {k: v.to("cuda") for k, v in inputs.items()}

        with torch.no_grad():
            final_ids = self.sum_model.generate(
                **inputs,
                num_beams=6,
                max_new_tokens=500,
                min_new_tokens=250,
                no_repeat_ngram_size=3,
                repetition_penalty=1.2,
                length_penalty=1.3,
            )

        final_summary = self.sum_tokenizer.decode(
            final_ids[0],
            skip_special_tokens=True
        )

        return final_summary.strip()

    # ==========================================================
    # TOPIC EXTRACTION (FROM SUMMARY)
    # ==========================================================

    def extract_topics(self, summary):

        print("--- Extracting Topics From Summary ---")

        sentences = re.split(r'(?<=[.!?])\s+', summary)

        topics = []
        seen = set()

        for s in sentences:
            s = s.strip()

            if len(s) > 50:
                short = s[:150]
                if short not in seen:
                    topics.append(short)
                    seen.add(short)

        return topics[:5]

    # ==========================================================
    # ACTION ITEM EXTRACTION (FILTERED + CLASSIFIED)
    # ==========================================================

    def rewrite_action_sentence(self, sentence: str) -> str:
        """
        Rewrite a raw transcript action sentence into a cleaner action-item line.
        This is ONLY used for LOW priority action items (no deadline).
        If generation fails, returns the original sentence.
        """
        s = re.sub(r"\s+", " ", sentence).strip()
        if len(s) < 25:
            return s

        # Light prompt to encourage "tas\\k-like" formation
        prompt = f"Rewrite as a clear action item: {s}"

        inputs = self.sum_tokenizer(
            prompt,
            return_tensors="pt",
            truncation=True,
            max_length=256
        )

        if torch.cuda.is_available():
            inputs = {k: v.to("cuda") for k, v in inputs.items()}

        try:
            self.sum_model.eval()
            with torch.no_grad():
                out = self.sum_model.generate(
                    **inputs,
                    num_beams=4,
                    max_new_tokens=40,
                    min_new_tokens=10,
                    no_repeat_ngram_size=3,
                    repetition_penalty=1.1,
                    length_penalty=1.0,
                    early_stopping=True
                )

            rewritten = self.sum_tokenizer.decode(out[0], skip_special_tokens=True).strip()

            # Safety checks
            if not rewritten or len(rewritten.split()) < 4:
                return s
            if rewritten.lower() == s.lower():
                return s

            # Nice formatting
            rewritten = rewritten[0].upper() + rewritten[1:]
            if not rewritten.endswith((".", "!", "?")):
                rewritten += "."

            return rewritten

        except Exception:
            return s

    def extract_action_items(self, text):
        print("--- Loading Custom AMI-Trained RoBERTa ---")

        model_path = r"D:\action_item_detection_model\content\drive\MyDrive\action_item_detection_model"

        classifier = pipeline(
            "text-classification",
            model=model_path,
            tokenizer=model_path,
            device=self.device
        )

        # -----------------------------
        # Helpers
        # -----------------------------
        def normalize(s: str) -> str:
            s = re.sub(r"\s+", " ", s).strip()
            s = re.sub(r"^(okay|so|well|now|right|yeah)\s*,?\s*", "", s, flags=re.I)
            return s

        NEGATIVE_PATTERNS = [
            r"\bcan we have your update\b",
            r"\bcould we have your update\b",
            r"\bcan we move on\b",
            r"\blet'?s move on\b",
            r"\bnext agenda item\b",
            r"\bany update\b",
            r"\bprovide an update\b",
            r"\bwhat'?s the update\b",
            r"\bthanks everyone\b",
            r"\bwelcome\b",
            r"\bmeeting minutes\b",
        ]

        POSITIVE_PATTERNS = [
            r"\bplease\s+(review|send|share|submit|update|finish|complete|prepare|organize|schedule|book|fix|implement|deploy|test|circulate|email|call)\b",
            r"\b(can|could)\s+you\s+(review|send|share|submit|update|finish|complete|prepare|organize|schedule|book|fix|implement|deploy|test|circulate|email|call)\b",
            r"\b(i|we)\s+will\s+(review|send|share|submit|update|finish|complete|prepare|organize|schedule|book|fix|implement|deploy|test|circulate|email|call)\b",
            r"\b(assign|action)\b",
            r"\bneed to\b",
            r"\bmake sure\b",
            r"\bensure\b",
        ]

        DEADLINE_PATTERNS = [
            (r"\b(today|tonight|eod|end of day|by end of day|by tonight)\b", "today/EOD"),
            (r"\b(tomorrow|tmrw)\b", "tomorrow"),
            (r"\b(this week|by this week|end of this week)\b", "this week"),
            (r"\b(next week|by next week|end of next week)\b", "next week"),
            (r"\b(monday|tuesday|wednesday|thursday|friday|saturday|sunday)\b", "weekday"),
            (r"\bby\s+\d{1,2}(:\d{2})?\s*(am|pm)\b", "time"),
            (r"\bby\s+\d{1,2}\s*(am|pm)\b", "time"),
        ]

        def extract_deadline(sentence: str):
            s = sentence.lower()
            for pat, label in DEADLINE_PATTERNS:
                if re.search(pat, s):
                    return label
            return None

        def is_negative(sentence: str) -> bool:
            s = sentence.lower()
            return any(re.search(p, s) for p in NEGATIVE_PATTERNS)

        def has_positive_signal(sentence: str) -> bool:
            s = sentence.lower()
            return any(re.search(p, s) for p in POSITIVE_PATTERNS)

        # -----------------------------
        # 1) Sentence split
        # -----------------------------
        sentences = re.split(r'(?<=[.!?])\s+', text)
        sentences = [normalize(s) for s in sentences if len(normalize(s)) >= 25]

        if not sentences:
            return []

        # -----------------------------
        # 2) Classifier FIRST
        # -----------------------------
        results = classifier(sentences)

        action_items = []
        seen = set()

        for sent, res in zip(sentences, results):

            if res["label"] != "LABEL_1" or res["score"] < 0.70:
                continue

            if is_negative(sent):
                continue

            if not has_positive_signal(sent):
                continue

            key = re.sub(r"[^a-z0-9]+", " ", sent.lower()).strip()
            if key in seen:
                continue
            seen.add(key)

            deadline = extract_deadline(sent)

            if deadline is not None:
                priority = "High"
                final_task = sent.strip()
            else:
                priority = "Low"
                final_task = self.rewrite_action_sentence(sent)

            action_items.append({
                "task": final_task,
                "priority": priority
            })

        del classifier
        import gc
        gc.collect()

        return action_items


def fact_check_document(text, api_key):

    client = OpenAI(api_key=api_key)
    
    prompt = f"""
    Perform a deep analysis of the provided text in two stages:

    Stage 1: FACT-CHECK
    Verify all names, dates, financial figures, and technical claims. Identify any internal contradictions or errors.

    Stage 2: VERIFIED SUMMARY
    Based on your analysis, provide exactly 10 comprehensive bullet points summarizing the most important aspects of the document. 
    - Ensure the points cover the beginning, middle, and end of the text.
    - If a fact was found to be incorrect in Stage 1, note the correction in the summary.

    Text: {text[:7000]}
    """
    
    try:
        response = client.chat.completions.create(
            model="gpt-4o", 
            messages=[{"role": "system", "content": "You are a professional fact-checker."},
                      {"role": "user", "content": prompt}],
            temperature=0
        )
        return response.choices[0].message.content
    except Exception as e:
        return f"API Error: {e}"

# In your main execution block:
if __name__ == "__main__":
    # Pull the key from .env safely
    MY_API_KEY = os.getenv("OPENAI_API_KEY")

# ==============================================================================
# 4. MAIN EXECUTION MENU
# ==============================================================================

def process_file(file_path, file_category, processor, api_key=None):
    if not os.path.exists(file_path):
        print(f"File not found: {file_path}")
        return

    # 1. Extraction
    if file_path.lower().endswith(('.mp4', '.avi', '.mov')):
        text = extract_text_from_video(file_path)
    elif file_path.lower().endswith(('.mp3', '.wav', '.mpeg')):
        text = extract_text_from_audio(file_path)
    elif file_path.lower().endswith(('.pdf', '.docx', '.pptx')):
        text = extract_text_from_document(file_path)
    else:
        print("Unsupported format.")
        return


    if not text:
        print("Could not extract any text from the file.")
        return
    

    # 2. NLP Processing
    print("\n--- Processing Results ---")
    user_query = "Summarize key decisions, updates, and next steps."  # you can customize this
    summary = processor.generate_summary(text)
    print(f"\n[SUMMARY]\n{summary}")

    if file_category in ["meeting", "lecture"]:
        topics = processor.extract_topics(summary)
        print("\n[TOPICS]:")
        if topics:
            for idx, topic in enumerate(topics, 1):
                print(f"{idx}. {topic}")
        else:
            print("- None")
        
    if file_category == "meeting":
        actions = processor.extract_action_items(text)
        high_priority = [item for item in actions if item['priority'] == 'High']
        low_priority  = [item for item in actions if item['priority'] == 'Low']

        print("\n[🔥 HIGH PRIORITY ACTION ITEMS]:")
        for item in high_priority: print(f"- {item['task']}")
        if not high_priority: print("- None")

        print("\n[🟢 LOW PRIORITY / ROUTINE]:")
        for item in low_priority: print(f"- {item['task']}")
        if not low_priority: print("- None")

    if file_category == "document" and api_key:
        print("\n[FACT CHECK]:")
        print(fact_check_document(text, api_key))

if __name__ == "__main__":
    try:
        print("==================================================")
        print("             NLP PROCESSING SYSTEM                ")
        print("==================================================")
        print("What would you like to do?")
        print("Process a single file (Video, Audio or Document)")

        # Initialize the NLP models
        processor = NLPProcessor()
        print("dtype:", next(processor.sum_model.parameters()).dtype)
        print("device:", "cuda" if torch.cuda.is_available() else "cpu")

        file_path = input("Please paste the path to your file: ").strip()
            # Remove any accidental quotation marks around the path
        file_path = file_path.strip('"').strip("'")
        file_category = input("What type of file is this? (meeting, lecture, document): ").strip().lower()
            
            # Optional: Add your OpenAI key if you want to use the Fact Checking feature
        MY_API_KEY = os.getenv("OPENAI_API_KEY") 
            
        print("\nStarting process... please wait.")
        process_file(file_path, file_category, processor, api_key=MY_API_KEY)
            

    except KeyboardInterrupt:
        print("\nProcess canceled by user.")